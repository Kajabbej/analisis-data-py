import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    
    # Mengeset ukuran slide menjadi widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Skema Warna
    c_navy = RGBColor(10, 37, 64)
    c_gray = RGBColor(80, 80, 80)
    c_light_gray = RGBColor(245, 245, 245)
    c_white = RGBColor(255, 255, 255)
    c_green = RGBColor(46, 117, 89)
    c_red = RGBColor(192, 0, 0)
    
    # Font Settings Helper
    def set_font(run, size_pt, bold=False, italic=False, color=c_gray, font_name="Calibri"):
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    # Helper untuk menambahkan header standar pada slide isi
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        set_font(p.runs[0], 28, bold=True, color=c_navy)
        return title_box

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Background)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank Layout
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Background Shape
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_navy
    bg.line.color.rgb = c_navy
    
    # Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "ANALISIS PERBANDINGAN METODE SAMPLING"
    set_font(p1.runs[0], 36, bold=True, color=c_white)
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Simple Random Sampling vs Stratified Random Sampling\nPada Data Pendapatan Harian E-Commerce"
    set_font(p2.runs[0], 20, bold=False, color=RGBColor(200, 220, 240))
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "\nDisusun oleh: Tim Kelompok Statistika"
    set_font(p3.runs[0], 14, italic=True, color=RGBColor(170, 190, 210))
    p3.alignment = PP_ALIGN.LEFT

    # ==========================================
    # SLIDE 2: LATAR BELAKANG & PEMBERSIHAN DATA
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide2, "Latar Belakang & Pembersihan Data")
    
    content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Sumber Data Utama", "Dataset transaksi Online Retail (Desember 2010 - Desember 2011) yang diagregasikan berdasarkan tanggal operasional harian."),
        ("Pembersihan Data Konsisten", "Seluruh baris transaksi bernilai negatif (refund/pembatalan) dihapus (filter Total Harga > 0) untuk menghindari bias rata-rata dan penyimpangan data."),
        ("Populasi Bersih (N)", "Menghasilkan total N = 304 hari operasional aktif yang siap dianalisis secara statistik."),
        ("Penggabungan Rentang Tahun", "Data tahun 2010 dan 2011 langsung digabungkan untuk mendapatkan jumlah hari per kategori secara riil tanpa sekat tahun."),
        ("Ekstraksi Kategori Hari", "Kolom nama hari (Senin, Selasa, Rabu, Kamis, Jumat, Minggu) diekstrak secara otomatis dari kolom tanggal asli.")
    ]
    
    for idx, (title, desc) in enumerate(bullets):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.space_after = Pt(12)
        
        r1 = p.add_run()
        r1.text = f"•  {title}: "
        set_font(r1, 16, bold=True, color=c_navy)
        
        r2 = p.add_run()
        r2.text = desc
        set_font(r2, 16, bold=False, color=c_gray)

    # ==========================================
    # SLIDE 3: METODOLOGI PENARIKAN SAMPEL (20%)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide3, "Metodologi Penarikan Sampel (Proporsi 20%)")
    
    # Kolom Kiri: SRS
    col1_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p_srs_title = tf1.paragraphs[0]
    p_srs_title.text = "1. Simple Random Sampling (SRS)"
    set_font(p_srs_title.runs[0], 20, bold=True, color=c_red)
    p_srs_title.space_after = Pt(10)
    
    srs_text = [
        "Sampel ditarik secara acak penuh dari keseluruhan 304 populasi tanpa mempertimbangkan kategori hari.",
        "Jumlah sampel yang terambil: n = 61 hari operasional.",
        "Setiap tanggal memiliki peluang terpilih yang sama.",
        "Kelemahan: Berisiko melewatkan kelompok hari tertentu atau nilai ekstrem (outlier) karena murni acak."
    ]
    for text in srs_text:
        p = tf1.add_paragraph()
        p.text = f"• {text}"
        set_font(p.runs[0], 15, color=c_gray)
        p.space_after = Pt(8)
        
    # Kolom Kanan: Stratified
    col2_box = slide3.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    
    p_strat_title = tf2.paragraphs[0]
    p_strat_title.text = "2. Stratified Random Sampling"
    set_font(p_strat_title.runs[0], 20, bold=True, color=c_green)
    p_strat_title.space_after = Pt(10)
    
    strat_text = [
        "Populasi dibagi menjadi 6 strata berdasarkan nama hari (Senin - Minggu).",
        "Sampel diambil tepat 20% secara acak dari masing-masing strata hari.",
        "Jumlah sampel yang terambil: n = 60 hari operasional.",
        "Kelebihan: Menjamin seluruh hari terwakili secara adil dan proporsional sesuai struktur populasi aslinya."
    ]
    for text in strat_text:
        p = tf2.add_paragraph()
        p.text = f"• {text}"
        set_font(p.runs[0], 15, color=c_gray)
        p.space_after = Pt(8)

    # ==========================================
    # SLIDE 4: HASIL STATISTIK DESKRIPTIF UTAMA (TABEL)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide4, "Tabel Perbandingan Parameter & Statistik Utama")
    
    # Tambah Tabel
    rows, cols = 7, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5)
    table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Parameter Statistik", "Data Populasi", "Sampel Simple Random (SRS)", "Sampel Stratified"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_navy
        set_font(cell.text_frame.paragraphs[0].runs[0], 15, bold=True, color=c_white)
        
    data = [
        ["Ukuran Data (N atau n)", "N = 304", "n = 61", "n = 60"],
        ["Rata-rata (Mean)", "32,070.11", "31,318.49", "31,456.21"],
        ["Standar Deviasi (SD)", "17,335.96", "16,910.17", "19,032.75"],
        ["Standard Error (SE)*", "-", "1,935.75", "2,226.79"],
        ["Nilai Minimum (Min)", "3,457.11", "4,137.62", "6,134.57"],
        ["Nilai Maksimum (Max)", "112,141.11", "75,244.43", "112,141.11"]
    ]
    
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            # Zebra striping
            if (r_idx + 1) % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = c_light_gray
            set_font(cell.text_frame.paragraphs[0].runs[0], 14, color=c_navy if c_idx == 0 else c_gray)

    # Footnote untuk FPC
    footnote_box = slide4.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.8))
    tf_fn = footnote_box.text_frame
    tf_fn.word_wrap = True
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "* Catatan: Nilai Standard Error (SE) telah disesuaikan menggunakan Finite Population Correction (FPC) karena fraksi sampling mencapai 20% (> 5% populasi). SE untuk Stratified dihitung berdasarkan rumus variansi tertimbang gabungan strata."
    set_font(p_fn.runs[0], 11, italic=True, color=c_gray)


    # ==========================================
    # SLIDE 5: TABEL STRATA HARI (STRATIFIED)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide5, "Perbandingan Parameter per Hari (Stratified)")
    
    # Tambah Tabel Strata
    rows, cols = 7, 7
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5)
    table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers_strata = ["Hari", "Pop N", "Samp n", "Pop Mean", "Samp Mean", "Pop SD", "Samp SD"]
    for i, h in enumerate(headers_strata):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_navy
        set_font(cell.text_frame.paragraphs[0].runs[0], 14, bold=True, color=c_white)
        
    data_strata = [
        ["Senin", "47", "9", "33,800.20", "37,327.17", "17,573.18", "29,219.61"],
        ["Selasa", "52", "10", "37,811.21", "33,851.66", "17,836.12", "18,028.97"],
        ["Rabu", "52", "10", "33,379.10", "32,312.10", "16,570.84", "17,541.98"],
        ["Kamis", "53", "11", "39,858.85", "36,130.45", "16,308.23", "17,260.91"],
        ["Jumat", "50", "10", "30,812.22", "28,169.15", "14,061.45", "18,169.50"],
        ["Minggu", "50", "10", "16,113.58", "21,066.42", "10,243.06", "11,025.08"]
    ]
    
    for r_idx, row in enumerate(data_strata):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            if (r_idx + 1) % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = c_light_gray
            # Highlight nama hari
            is_bold = (c_idx == 0)
            set_font(cell.text_frame.paragraphs[0].runs[0], 13, bold=is_bold, color=c_navy if c_idx == 0 else c_gray)

    # ==========================================
    # SLIDE 6: KESIMPULAN ANALISIS DISTRIBUSI KDE
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide6, "Analisis Keterwakilan & Grafik Distribusi (KDE)")
    
    content_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        ("Representasi Rata-rata (Mean)", "Kedua metode sampling sukses memperkirakan rata-rata populasi dengan selisih yang sangat tipis (di bawah 2%). Rata-rata sampel Stratified (31,456.21) sedikit lebih dekat ke populasi dibanding SRS (31,318.49)."),
        ("Penangkapan Data Ekstrem (Outlier)", "Sampel Stratified berhasil menangkap nilai transaksi maksimum populasi (112,141.11). Sebaliknya, sampel SRS melewatkannya (maksimum sampel SRS hanya 75,244.43). Ini membuktikan Stratified lebih unggul dalam menjaga variabilitas sebaran data."),
        ("Penyimpangan Data (SD)", "SRS menghasilkan SD yang lebih kecil karena hilangnya outlier. Sedangkan Stratified menghasilkan SD yang sedikit lebih besar akibat sensitivitas ukuran sampel yang kecil (n sekitar 10 hari) per kategori strata hari."),
        ("Bentuk Kurva KDE", "Grafik KDE membuktikan kurva Stratified (hijau putus-putus) memiliki kelengkungan yang lebih presisi menempel pada kurva populasi asli dibandingkan dengan kurva SRS (merah putus-putus).")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.space_after = Pt(15)
        
        r1 = p.add_run()
        r1.text = f"✔  {title}: "
        set_font(r1, 16, bold=True, color=c_navy)
        
        r2 = p.add_run()
        r2.text = desc
        set_font(r2, 16, bold=False, color=c_gray)

    # ==========================================
    # SLIDE 7: RENCANA PENGUJIAN HIPOTESIS LANJUTAN
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide7, "Rencana Pengujian Hipotesis Kelompok")
    
    # Kolom Kiri: Uji-t
    col1_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p_t_title = tf1.paragraphs[0]
    p_t_title.text = "1. Uji Keterwakilan (One-Sample t-Test)"
    set_font(p_t_title.runs[0], 18, bold=True, color=c_navy)
    p_t_title.space_after = Pt(10)
    
    t_points = [
        "Tujuan: Menguji secara formal apakah rata-rata sampel berbeda secara signifikan dengan rata-rata populasi.",
        "Hipotesis Nol (H0): Rata-rata sampel = Rata-rata populasi (Sampel mewakili populasi).",
        "Hipotesis Alternatif (H1): Rata-rata sampel ≠ Rata-rata populasi (Sampel tidak mewakili populasi).",
        "Rumus t-hitung: t = (Rata-rata Sampel - Rata-rata Populasi) / SE_sampel (dengan FPC)"
    ]
    for text in t_points:
        p = tf1.add_paragraph()
        p.text = f"• {text}"
        set_font(p.runs[0], 14, color=c_gray)
        p.space_after = Pt(8)
        
    # Kolom Kanan: ANOVA
    col2_box = slide7.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    
    p_anova_title = tf2.paragraphs[0]
    p_anova_title.text = "2. Uji Beda Kelompok (ANOVA)"
    set_font(p_anova_title.runs[0], 18, bold=True, color=c_navy)
    p_anova_title.space_after = Pt(10)
    
    anova_points = [
        "Tujuan: Membuktikan apakah ada perbedaan pendapatan harian yang signifikan antar hari dalam seminggu.",
        "Hipotesis Nol (H0): Rata-rata pendapatan Senin = Selasa = Rabu = Kamis = Jumat = Minggu.",
        "Hipotesis Alternatif (H1): Minimal ada satu hari yang memiliki rata-rata pendapatan berbeda secara signifikan.",
        "Analisis dapat dilakukan di Excel menggunakan menu 'Data Analysis -> Anova: Single Factor' atau menggunakan Python scipy.stats."
    ]
    for text in anova_points:
        p = tf2.add_paragraph()
        p.text = f"• {text}"
        set_font(p.runs[0], 14, color=c_gray)
        p.space_after = Pt(8)

    # Save Presentation
    output_dir = "laporan"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    output_path = os.path.join(output_dir, "presentasi_analisis_sampling.pptx")
    prs.save(output_path)
    print(f"[OK] Berhasil membuat berkas presentasi: {output_path}")

if __name__ == '__main__':
    create_presentation()
